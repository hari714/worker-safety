"""
Generate Third Review Documents (Topics 4-8) for Workers Safety System
Arasu Engineering College, Kumbakonam - CSE A Final Year
"""

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ============================================================
# CONTENT DATA
# ============================================================

PROJECT_TITLE = "AI-Powered Workers Safety Monitoring System"
COLLEGE = "Arasu Engineering College, Kumbakonam - 612 501"
DEPT = "Department of Computer Science and Engineering"
ACADEMIC_YEAR = "Academic Year 2025-2026 (Even Sem)"
REVIEW = "Third Review"

# Topic 4: System Requirements
HARDWARE_REQS = [
    ("Processor", "Intel Core i5 (8th Gen or above) / AMD Ryzen 5 or equivalent"),
    ("RAM", "8 GB minimum (16 GB recommended for model training)"),
    ("GPU", "NVIDIA GPU with CUDA support (optional, recommended for real-time detection)"),
    ("Storage", "10 GB+ free disk space (for models, datasets, and logs)"),
    ("Camera", "USB Webcam / IP Camera with RTSP stream support"),
    ("Display", "1280 x 720 minimum resolution"),
    ("Network", "Internet connection for email notifications"),
]

SOFTWARE_REQS = [
    ("Operating System", "Windows 10/11, Ubuntu 20.04+, or macOS 12+"),
    ("Programming Language", "Python 3.8 or above"),
    ("Web Framework", "Flask 2.3.0"),
    ("Database", "SQLite 3 (built-in with Python)"),
    ("ORM", "SQLAlchemy 2.0+"),
    ("AI/ML Framework", "PyTorch 2.0, Ultralytics YOLOv8"),
    ("Computer Vision", "OpenCV 4.8.0"),
    ("Face Detection", "MTCNN (Multi-task Cascaded CNN)"),
    ("Face Recognition", "FaceNet (512-dim embeddings)"),
    ("Image Processing", "Pillow 10.0, NumPy 1.24"),
    ("IDE", "VS Code / PyCharm"),
    ("Version Control", "Git"),
    ("Browser", "Chrome / Firefox / Edge (latest)"),
]

# Topic 5: System Design
SYSTEM_DESIGN_DESC = (
    "The Workers Safety Monitoring System follows a three-layer architecture: "
    "Input Layer, Processing Layer, and Action Layer. The system captures live "
    "video feeds from cameras, processes each frame through AI models for PPE "
    "detection and face recognition, and takes appropriate actions such as "
    "logging violations, sending email alerts, and updating the dashboard."
)

ARCHITECTURE_LAYERS = [
    ("Input Layer", [
        "USB Webcam / IP Camera capture",
        "Video file upload",
        "Image upload via web interface",
        "Frame extraction at 30 FPS",
    ]),
    ("Processing Layer", [
        "Frame preprocessing (resize to 640x640, normalize, color conversion)",
        "YOLOv8 PPE Detection - detects Helmet, Gloves, Vest, Boots, Goggles",
        "MTCNN Face Detection - locates faces in frames",
        "FaceNet Face Recognition - generates 512-dim embeddings for identification",
        "Compliance Checker - compares detected PPE against required PPE",
    ]),
    ("Action Layer", [
        "SQLite Database - stores worker info, violations, compliance history",
        "Email Notification System - sends alerts via SMTP for violations",
        "Web Dashboard - displays real-time detection results",
        "REST API - provides endpoints for integration",
        "Violation Logger - captures screenshots and logs non-compliance",
    ]),
]

DATA_FLOW_STEPS = [
    "Camera/video input is captured and frames are extracted",
    "Each frame is preprocessed (resized, normalized, color-converted)",
    "Parallel AI processing: PPE detection (YOLOv8) + Face recognition (MTCNN/FaceNet)",
    "Compliance check: detected PPE compared against safety requirements",
    "If violation found: log to database, send email alert, display on dashboard",
    "If compliant: log compliant status to history",
    "Results displayed on web UI with annotated bounding boxes and confidence scores",
]

# Topic 6: User Interface
UI_PAGES = [
    ("Home / Upload Page", [
        "Clean black-and-white minimal design",
        "Project title header with safety icon",
        "Drag-and-drop image upload area with visual feedback",
        "File input selector with image preview",
        "Submit button for PPE analysis",
        "Responsive layout for desktop and mobile",
    ]),
    ("Detection Results Page", [
        "Annotated image with color-coded bounding boxes around detected PPE",
        "Overall compliance status indicator (PASS in green / FAIL in red)",
        "PPE Checklist showing each item with detection confidence percentage",
        "Statistics panel: Detected count, Missing count, Total objects",
        "Color coding: Green = Found, Red = Missing, Yellow = Low confidence",
        "Option to upload another image for analysis",
    ]),
    ("Live Monitoring View (OpenCV Window)", [
        "Real-time video feed with PPE detection overlays",
        "Bounding boxes with class labels and confidence scores",
        "FPS counter for performance monitoring",
        "Keyboard controls: 'q' to quit, 's' to save screenshot",
    ]),
]

# Topic 7: List of Modules
MODULES = [
    ("PPE Detection Module",
     "Core AI module using YOLOv8 for detecting 5 types of Personal Protective Equipment "
     "(Helmet, Gloves, Vest, Boots, Goggles) in real-time with 95%+ accuracy.",
     "Implemented",
     ["yolo_model.py - YOLOv8 model training with data augmentation",
      "ppe_classes.py - PPE data classes and type definitions",
      "inference.py - Real-time PPE detector and compliance checker"]),

    ("Face Recognition Module",
     "Worker identification system using MTCNN for face detection and FaceNet for "
     "generating 512-dimensional face embeddings for matching.",
     "Skeleton",
     ["face_detection.py - MTCNN-based face detection",
      "face_embedding.py - FaceNet embedding generator",
      "worker_identification.py - Cosine similarity matching",
      "face_database.py - Store/load face embeddings (.pkl files)"]),

    ("Worker Management Module",
     "Database operations for managing worker records including registration, "
     "updates, and queries using SQLAlchemy ORM.",
     "Skeleton",
     ["models.py - SQLAlchemy database models (Workers, Violations, Compliance)",
      "db_operations.py - CRUD operations for all tables",
      "worker_registry.py - Worker registration and face enrollment"]),

    ("Notification System Module",
     "Automated email alert system that sends violation notifications to supervisors "
     "with details of missing PPE and captured screenshots.",
     "Skeleton",
     ["email_sender.py - SMTP email service",
      "email_templates.py - HTML email templates with violation details",
      "notification_queue.py - Queue management for batch notifications"]),

    ("Monitoring Module",
     "Main processing loop that integrates all modules for continuous video "
     "monitoring, compliance checking, and violation logging.",
     "Skeleton",
     ["video_processor.py - Main video processing loop",
      "compliance_checker.py - PPE compliance verification",
      "violation_logger.py - Screenshot capture and violation database logging"]),

    ("REST API Module",
     "Flask-based REST API providing HTTP endpoints for image upload, worker "
     "management, violation logs, and compliance reports.",
     "Skeleton",
     ["app.py - Flask application factory",
      "routes.py - API endpoint definitions",
      "middleware.py - Authentication and request validation"]),

    ("Web Application Module",
     "Flask web interface with drag-and-drop image upload, real-time PPE detection "
     "visualization, and compliance status display.",
     "Implemented",
     ["web_app.py - Complete Flask web app with HTML/CSS/JS UI",
      "Drag-and-drop upload with image preview",
      "Annotated result display with color-coded bounding boxes"]),

    ("Configuration & Utilities Module",
     "System configuration management, logging, input validation, and helper "
     "functions used across all modules.",
     "Implemented",
     ["config/settings.py - Environment variable loader",
      "config/database.py - SQLAlchemy database setup",
      "config/email_config.py - SMTP configuration",
      "utils/logger.py - File and console logging",
      "utils/validators.py - Email, phone, worker ID validation",
      "utils/helpers.py - Timestamp and filename utilities"]),
]

# Topic 8: Backend Processes of Completed Modules
BACKEND_PROCESSES = [
    ("PPE Detection Backend Process", [
        ("Model Loading", "The trained YOLOv8 model (ppe_detection_best.pt, 6.2 MB) is loaded into memory. "
         "The model supports GPU acceleration via CUDA if available, otherwise falls back to CPU."),
        ("Image Preprocessing", "Input images are resized to 640x640 pixels, normalized, and converted "
         "from BGR to RGB color space using OpenCV for model compatibility."),
        ("Inference Pipeline", "The YOLOv8 model processes the image and returns detection results including "
         "bounding box coordinates, class labels (0-4 for 5 PPE types), and confidence scores."),
        ("Confidence Filtering", "Detections with confidence below the threshold (0.5) are filtered out. "
         "Additionally, detections covering more than 30% of the image area are considered false positives and removed."),
        ("Compliance Check", "Detected PPE items are compared against the full set of 5 required items. "
         "Missing items are identified and reported with their status."),
        ("Result Annotation", "Bounding boxes are drawn on the image with color coding: Green (Helmet), "
         "Blue (Gloves), Magenta (Vest), Red (Boots), Cyan (Goggles). Labels include class name and confidence %."),
    ]),
    ("Web Application Backend Process", [
        ("Flask Server Initialization", "The Flask app starts on port 5000, configures maximum upload size (16 MB), "
         "and loads the PPE detection model into memory."),
        ("Request Handling", "When a POST request is received with an image file, the server validates the file type "
         "(JPEG, PNG, GIF, BMP, WEBP) and reads it into a NumPy array."),
        ("Image Processing Pipeline", "The uploaded image is decoded using OpenCV, converted from BGRA to BGR if needed, "
         "and passed to the PPE detector for inference."),
        ("Result Generation", "Detection results are compiled into a structured format with PPE checklist, "
         "compliance status, and statistics. The annotated image is encoded to base64 for HTML embedding."),
        ("Template Rendering", "Jinja2 templates render the results page with the annotated image, PPE checklist "
         "with confidence values, and overall compliance status (PASS/FAIL)."),
    ]),
    ("Model Training Backend Process", [
        ("Dataset Preparation", "The training dataset contains 3,810 images with 19,848 labeled PPE objects "
         "split into train (2,667), validation (571), and test (572) sets in YOLO format."),
        ("Data Augmentation", "Training applies augmentations: HSV shift (h=0.015, s=0.7, v=0.4), "
         "rotation (10 deg), translation (10%), scaling (50%), horizontal flip (50%), and mosaic."),
        ("Training Configuration", "YOLOv8 medium model trained for 100 epochs with batch size 32, "
         "SGD optimizer (momentum=0.937), initial LR=0.01, early stopping patience=20."),
        ("Checkpoint Management", "Model checkpoints saved every 5 epochs. Best model (by validation mAP) "
         "is automatically copied to models/ppe_detection_best.pt."),
        ("Performance Metrics", "Training logs mAP, precision, recall, and loss curves. "
         "Final model achieves 95%+ detection accuracy across all 5 PPE classes."),
    ]),
    ("Configuration & Logging Backend Process", [
        ("Environment Loading", "Settings loaded from .env file using python-dotenv. Includes confidence thresholds, "
         "camera settings, database path, and email credentials."),
        ("Database Connection", "SQLAlchemy creates a thread-safe SQLite connection with session management. "
         "Database file stored at safety_db.sqlite in project root."),
        ("Logging System", "Dual-output logging: console (INFO level) and file (DEBUG level) at logs/safety_system.log. "
         "Includes timestamps, module names, and log levels."),
        ("Input Validation", "Validators ensure data integrity: email format validation, worker ID format check, "
         "phone number validation, and file type verification."),
    ]),
]


# ============================================================
# PDF GENERATION
# ============================================================

class ReviewPDF(FPDF):
    def __init__(self):
        super().__init__()
        self.set_auto_page_break(auto=True, margin=25)

    def header(self):
        pass

    def footer(self):
        pass

    def section_title(self, num, title):
        self.add_page()
        self.set_font("Helvetica", "B", 18)
        self.set_text_color(0, 0, 0)
        self.cell(0, 14, f"{num}. {title}", new_x="LMARGIN", new_y="NEXT")
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(6)

    def sub_heading(self, title):
        self.ln(4)
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(30, 30, 30)
        self.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT")
        self.ln(2)

    def body_text(self, text):
        self.set_font("Helvetica", "", 11)
        self.set_text_color(40, 40, 40)
        self.multi_cell(0, 6.5, text)
        self.ln(2)

    def bullet(self, text, indent=15):
        self.set_font("Helvetica", "", 10.5)
        self.set_text_color(40, 40, 40)
        x = self.get_x()
        self.cell(indent, 6, "")
        # bullet character
        self.set_font("Helvetica", "B", 10.5)
        self.cell(5, 6, ">")
        self.set_font("Helvetica", "", 10.5)
        self.multi_cell(0, 6, f" {text}")
        self.ln(1)

    def table_row(self, col1, col2, widths=(55, 125), bold_first=True):
        self.set_font("Helvetica", "B" if bold_first else "", 10.5)
        self.set_text_color(40, 40, 40)
        x = self.get_x()
        y = self.get_y()
        self.cell(widths[0], 7, f"  {col1}")
        self.set_font("Helvetica", "", 10.5)
        self.multi_cell(widths[1], 7, col2)
        self.ln(1)


def generate_pdf():
    pdf = ReviewPDF()

    # ---- Topic 4: System Requirements ----
    pdf.section_title(4, "System Requirements")
    pdf.body_text(
        "This section outlines the hardware and software requirements needed to deploy and run "
        "the AI-Powered Workers Safety Monitoring System effectively."
    )

    pdf.sub_heading("4.1 Hardware Requirements")
    for name, desc in HARDWARE_REQS:
        pdf.table_row(name, desc)

    pdf.sub_heading("4.2 Software Requirements")
    for name, desc in SOFTWARE_REQS:
        pdf.table_row(name, desc)

    # ---- Topic 5: System Design ----
    pdf.section_title(5, "System Design")
    pdf.body_text(SYSTEM_DESIGN_DESC)

    pdf.sub_heading("5.1 Three-Layer Architecture")
    for layer_name, items in ARCHITECTURE_LAYERS:
        pdf.ln(2)
        pdf.set_font("Helvetica", "B", 11.5)
        pdf.set_text_color(20, 20, 20)
        pdf.cell(0, 8, f"  {layer_name}", new_x="LMARGIN", new_y="NEXT")
        for item in items:
            pdf.bullet(item)

    pdf.sub_heading("5.2 Data Flow")
    for i, step in enumerate(DATA_FLOW_STEPS, 1):
        pdf.set_font("Helvetica", "", 10.5)
        pdf.set_text_color(40, 40, 40)
        pdf.cell(15, 7, "")
        pdf.set_font("Helvetica", "B", 10.5)
        pdf.cell(8, 7, f"{i}.")
        pdf.set_font("Helvetica", "", 10.5)
        pdf.multi_cell(0, 7, step)
        pdf.ln(1)

    pdf.sub_heading("5.3 Key Configuration Parameters")
    configs = [
        ("PPE Confidence", "0.5 (minimum detection threshold)"),
        ("Face Confidence", "0.6 (minimum face detection threshold)"),
        ("Similarity Threshold", "0.6 (cosine similarity for face matching)"),
        ("Image Size", "640 x 640 pixels (model input)"),
        ("Video FPS", "30 frames per second"),
        ("Frame Resolution", "1280 x 720 pixels"),
    ]
    for name, val in configs:
        pdf.table_row(name, val)

    # ---- Topic 6: User Interface ----
    pdf.section_title(6, "User Interface")
    pdf.body_text(
        "The system provides a web-based user interface built with Flask, featuring a modern "
        "minimal black-and-white design. The interface enables users to upload images for PPE "
        "detection analysis and view detailed results with visual annotations."
    )

    for page_name, features in UI_PAGES:
        pdf.sub_heading(f"  {page_name}")
        for f in features:
            pdf.bullet(f)

    # ---- Topic 7: List of Modules ----
    pdf.section_title(7, "List of Modules")
    pdf.body_text(
        "The system is organized into 8 modular components, each responsible for a specific "
        "functionality. This modular design ensures maintainability, scalability, and clear "
        "separation of concerns."
    )

    for i, (name, desc, status, files) in enumerate(MODULES, 1):
        pdf.sub_heading(f"7.{i} {name}  [{status}]")
        pdf.body_text(desc)
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(15, 7, "")
        pdf.cell(0, 7, "Key Files:", new_x="LMARGIN", new_y="NEXT")
        for f in files:
            pdf.bullet(f, indent=20)

    # ---- Topic 8: Backend Processes ----
    pdf.section_title(8, "Backend Processes of Completed Modules")
    pdf.body_text(
        "This section describes the backend processes and workflows of the modules that have "
        "been fully implemented and are operational in the current version of the system."
    )

    for proc_name, steps in BACKEND_PROCESSES:
        pdf.sub_heading(f"  {proc_name}")
        for step_name, step_desc in steps:
            pdf.set_font("Helvetica", "B", 10.5)
            pdf.set_text_color(30, 30, 30)
            pdf.cell(15, 7, "")
            pdf.cell(0, 7, step_name, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(50, 50, 50)
            pdf.cell(20, 6, "")
            pdf.multi_cell(170, 6, step_desc)
            pdf.ln(2)

    filepath = os.path.join(OUTPUT_DIR, "Third_Review_Topics_4_to_8.pdf")
    pdf.output(filepath)
    print(f"PDF saved: {filepath}")
    return filepath


# ============================================================
# PPT GENERATION
# ============================================================

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # dark navy

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 210, 230)
    p2.alignment = PP_ALIGN.CENTER

    # College info
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.4), Inches(1.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = COLLEGE
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(160, 175, 200)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf3.add_paragraph()
    p4.text = DEPT
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(140, 155, 180)
    p4.alignment = PP_ALIGN.CENTER
    p5 = tf3.add_paragraph()
    p5.text = ACADEMIC_YEAR
    p5.font.size = Pt(12)
    p5.font.color.rgb = RGBColor(140, 155, 180)
    p5.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, section_num, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 41, 59)

    # Number
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Topic {section_num}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(100, 180, 255)
    p.alignment = PP_ALIGN.CENTER

    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(34)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle if provided
    y_start = 1.3
    if subtitle:
        txBox_sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.5))
        tf_sub = txBox_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.italic = True
        p_sub.font.color.rgb = RGBColor(100, 100, 120)
        y_start = 1.7

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(y_start), Inches(8.8), Inches(5.5 - y_start + 1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 50)
        p.space_after = Pt(8)
        p.level = 0
        # Add bullet character
        p.text = f"  {bullet}"


def add_table_slide(prs, title, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Table
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 2

    if col_widths is None:
        col_widths = [Inches(3), Inches(5.8)]

    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3),
                                          sum(col_widths), Inches(0.35 * num_rows))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    paragraph.font.color.rgb = RGBColor(40, 40, 50)

            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 243, 248)


def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Title Slide ----
    add_title_slide(prs, PROJECT_TITLE, f"{REVIEW} - Topics 4 to 8")

    # ---- Topic 4: System Requirements ----
    add_section_slide(prs, 4, "System Requirements")

    hw_rows = [("Component", "Specification")] + [(n, d) for n, d in HARDWARE_REQS]
    add_table_slide(prs, "4.1 Hardware Requirements", hw_rows)

    # Split software into two slides
    sw1 = SOFTWARE_REQS[:7]
    sw2 = SOFTWARE_REQS[7:]
    sw_rows1 = [("Component", "Specification")] + [(n, d) for n, d in sw1]
    add_table_slide(prs, "4.2 Software Requirements (1/2)", sw_rows1)

    sw_rows2 = [("Component", "Specification")] + [(n, d) for n, d in sw2]
    add_table_slide(prs, "4.2 Software Requirements (2/2)", sw_rows2)

    # ---- Topic 5: System Design ----
    add_section_slide(prs, 5, "System Design")

    add_content_slide(prs, "5.1 System Overview", [
        "Three-layer architecture: Input, Processing, Action",
        "Real-time video processing with AI-powered detection",
        "YOLOv8 for PPE detection (5 equipment types)",
        "MTCNN + FaceNet for worker identification",
        "Automated alerts and compliance tracking",
    ])

    for layer_name, items in ARCHITECTURE_LAYERS:
        add_content_slide(prs, f"5.2 {layer_name}", items)

    add_content_slide(prs, "5.3 Data Flow", DATA_FLOW_STEPS)

    config_bullets = [f"{n}: {v}" for n, v in [
        ("PPE Confidence", "0.5 (minimum detection threshold)"),
        ("Face Confidence", "0.6 (minimum face detection threshold)"),
        ("Similarity Threshold", "0.6 (cosine similarity for face matching)"),
        ("Image Size", "640x640 pixels (model input)"),
        ("Video FPS", "30 frames per second"),
        ("Frame Resolution", "1280x720 pixels"),
    ]]
    add_content_slide(prs, "5.4 Key Configuration Parameters", config_bullets)

    # ---- Topic 6: User Interface ----
    add_section_slide(prs, 6, "User Interface")

    for page_name, features in UI_PAGES:
        add_content_slide(prs, f"6. {page_name}", features)

    # ---- Topic 7: List of Modules ----
    add_section_slide(prs, 7, "List of Modules")

    # Overview slide
    mod_overview = [("Module", "Status")] + [(m[0], m[2]) for m in MODULES]
    add_table_slide(prs, "7.1 Module Overview", mod_overview)

    # Individual module slides
    for i, (name, desc, status, files) in enumerate(MODULES, 1):
        bullets = [desc, ""] + [f"  {f}" for f in files]
        add_content_slide(prs, f"7.{i} {name}", [desc] + files, subtitle=f"Status: {status}")

    # ---- Topic 8: Backend Processes ----
    add_section_slide(prs, 8, "Backend Processes of Completed Modules")

    for proc_name, steps in BACKEND_PROCESSES:
        # Each process gets a slide
        bullets = [f"{step_name}: {step_desc[:90]}..." if len(step_desc) > 93
                   else f"{step_name}: {step_desc}"
                   for step_name, step_desc in steps[:5]]
        add_content_slide(prs, proc_name, bullets)

    # ---- Thank You Slide ----
    add_title_slide(prs, "Thank You", "Questions & Discussion")

    filepath = os.path.join(OUTPUT_DIR, "Third_Review_Topics_4_to_8.pptx")
    prs.save(filepath)
    print(f"PPTX saved: {filepath}")
    return filepath


# ============================================================
# MAIN
# ============================================================

if __name__ == "__main__":
    print("Generating Third Review Documents (Topics 4-8)...")
    print("=" * 50)
    pdf_path = generate_pdf()
    pptx_path = generate_pptx()
    print("=" * 50)
    print("Done! Files generated:")
    print(f"  PDF:  {pdf_path}")
    print(f"  PPTX: {pptx_path}")
